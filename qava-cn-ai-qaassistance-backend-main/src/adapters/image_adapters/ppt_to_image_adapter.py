from pptx import Presentation
from PIL import Image, ImageDraw

class PptToImageAdapter:
    def __init__(self, file_path):
        self.file_path = file_path

    def convert_to_image(self):
        presentation = Presentation(self.file_path)
        images = []
        
        for slide in presentation.slides:
            img = Image.new('RGB', (1280, 720), color = (255, 255, 255))
            d = ImageDraw.Draw(img)
            slide_text = "\n".join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
            d.text((10, 10), slide_text, fill=(0, 0, 0))
            images.append(img)

        return images
